from __future__ import annotations

import csv
import html
import hashlib
import io
import json
import os
import re
import urllib.error
import urllib.parse
import urllib.request
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

import numpy as np

try:
    import PyPDF2
except ImportError:
    PyPDF2 = None  # type: ignore[assignment]

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None  # type: ignore[assignment]

try:
    from openpyxl import load_workbook
except Exception:
    load_workbook = None  # type: ignore[assignment]

try:
    from pptx import Presentation
except Exception:
    Presentation = None  # type: ignore[assignment]

try:
    import xlrd
except Exception:
    xlrd = None  # type: ignore[assignment]


TEXT_LIKE_EXTENSIONS = {
    ".c",
    ".cpp",
    ".cs",
    ".css",
    ".go",
    ".java",
    ".js",
    ".jsx",
    ".kt",
    ".log",
    ".md",
    ".php",
    ".py",
    ".rb",
    ".rs",
    ".sql",
    ".swift",
    ".ts",
    ".tsx",
    ".txt",
    ".xml",
    ".yaml",
    ".yml",
}

IMAGE_EXTENSIONS = {".avif", ".gif", ".heic", ".heif", ".jpeg", ".jpg", ".png", ".webp"}


STOPWORDS = {
    "a",
    "about",
    "an",
    "and",
    "are",
    "as",
    "at",
    "be",
    "by",
    "for",
    "from",
    "how",
    "i",
    "in",
    "is",
    "it",
    "me",
    "my",
    "of",
    "on",
    "or",
    "please",
    "tell",
    "the",
    "to",
    "we",
    "what",
    "when",
    "where",
    "which",
    "who",
    "with",
    "you",
    "এবং",
    "একটি",
    "একটা",
    "এর",
    "কি",
    "কী",
    "কিভাবে",
    "কীভাবে",
    "আমি",
    "আমার",
    "জানতে",
    "চাই",
    "বলুন",
}

EMBEDDING_MODEL = "text-embedding-004"
EMBEDDING_DIMENSION = 256


@dataclass
class DocumentChunk:
    id: str
    source: str
    title: str
    text: str
    normalized_text: str
    tokens: set[str]
    embedding: list[float] = field(default_factory=list)


@dataclass
class RetrievalResult:
    chunk: DocumentChunk
    score: float


@dataclass
class RAGResponse:
    question: str
    answer: str
    category: str
    source_url: str | None
    language: str
    confidence: float
    found_match: bool
    used_gemini: bool
    matched_faq_id: str | None
    source_urls: list[str]
    source_titles: list[str]


def parse_pdf(file_bytes: bytes, filename: str) -> str:
    if PyPDF2 is None:
        raise ImportError("PyPDF2 is required for PDF parsing. Install it with pip install PyPDF2.")
    reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
    pages: list[str] = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            pages.append(text.strip())
    return "\n\n".join(pages)


def parse_text_file(file_bytes: bytes, filename: str) -> str:
    return file_bytes.decode("utf-8", errors="ignore")


def parse_csv_file(file_bytes: bytes, filename: str) -> str:
    rows: list[str] = []
    decoded = io.StringIO(file_bytes.decode("utf-8", errors="ignore"))
    delimiter = "\t" if Path(filename).suffix.lower() == ".tsv" else ","
    for row in csv.reader(decoded, delimiter=delimiter):
        cleaned = [cell.strip() for cell in row if cell.strip()]
        if cleaned:
            rows.append(" | ".join(cleaned))
    return "\n".join(rows)


def parse_docx_file(file_bytes: bytes, filename: str) -> str:
    if DocxDocument is None:
        raise ImportError("python-docx is required for DOCX parsing. Install it with pip install python-docx.")
    document = DocxDocument(io.BytesIO(file_bytes))
    paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
    table_rows: list[str] = []
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                table_rows.append(" | ".join(cells))
    return "\n\n".join([*paragraphs, *table_rows])


def parse_json_file(file_bytes: bytes, filename: str) -> str:
    decoded = file_bytes.decode("utf-8", errors="ignore")
    try:
        parsed = json.loads(decoded)
    except json.JSONDecodeError:
        return decoded
    return json.dumps(parsed, ensure_ascii=False, indent=2)


def parse_html_file(file_bytes: bytes, filename: str) -> str:
    decoded = file_bytes.decode("utf-8", errors="ignore")
    decoded = re.sub(r"(?is)<(script|style).*?>.*?</\1>", " ", decoded)
    decoded = re.sub(r"(?s)<[^>]+>", " ", decoded)
    return html.unescape(re.sub(r"\s+", " ", decoded)).strip()


def parse_rtf_file(file_bytes: bytes, filename: str) -> str:
    decoded = file_bytes.decode("utf-8", errors="ignore")
    decoded = re.sub(r"\\'[0-9a-fA-F]{2}", " ", decoded)
    decoded = re.sub(r"\\[a-zA-Z]+\d* ?", " ", decoded)
    decoded = decoded.replace("\\{", "{").replace("\\}", "}").replace("\\\\", "\\")
    decoded = re.sub(r"[{}]", " ", decoded)
    return re.sub(r"\s+", " ", decoded).strip()


def parse_xlsx_file(file_bytes: bytes, filename: str) -> str:
    if load_workbook is None:
        raise ImportError("openpyxl is required for XLSX parsing. Install it with pip install openpyxl.")

    workbook = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    sections: list[str] = []
    for sheet in workbook.worksheets:
        rows: list[str] = []
        for row in sheet.iter_rows(values_only=True):
            cleaned = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
            if cleaned:
                rows.append(" | ".join(cleaned))
        if rows:
            sections.append(f"Sheet: {sheet.title}\n" + "\n".join(rows))
    return "\n\n".join(sections)


def parse_xls_file(file_bytes: bytes, filename: str) -> str:
    if xlrd is None:
        raise ImportError("xlrd is required for legacy XLS parsing. Install it with pip install xlrd.")

    workbook = xlrd.open_workbook(file_contents=file_bytes)
    sections: list[str] = []
    for sheet in workbook.sheets():
        rows: list[str] = []
        for row_index in range(sheet.nrows):
            cleaned = [
                str(sheet.cell_value(row_index, column_index)).strip()
                for column_index in range(sheet.ncols)
                if str(sheet.cell_value(row_index, column_index)).strip()
            ]
            if cleaned:
                rows.append(" | ".join(cleaned))
        if rows:
            sections.append(f"Sheet: {sheet.name}\n" + "\n".join(rows))
    return "\n\n".join(sections)


def parse_pptx_file(file_bytes: bytes, filename: str) -> str:
    if Presentation is None:
        raise ImportError("python-pptx is required for PPTX parsing. Install it with pip install python-pptx.")

    presentation = Presentation(io.BytesIO(file_bytes))
    slides: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                parts.append(text.strip())
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        parts.append(" | ".join(cells))
        if parts:
            slides.append(f"Slide {index}\n" + "\n".join(parts))
    return "\n\n".join(slides)


def parse_uploaded_file(file_bytes: bytes, filename: str) -> str:
    ext = Path(filename).suffix.lower()
    if ext == ".pdf":
        return parse_pdf(file_bytes, filename)
    if ext in TEXT_LIKE_EXTENSIONS:
        return parse_text_file(file_bytes, filename)
    if ext in {".csv", ".tsv"}:
        return parse_csv_file(file_bytes, filename)
    if ext == ".docx":
        return parse_docx_file(file_bytes, filename)
    if ext == ".json":
        return parse_json_file(file_bytes, filename)
    if ext in {".html", ".htm"}:
        return parse_html_file(file_bytes, filename)
    if ext == ".rtf":
        return parse_rtf_file(file_bytes, filename)
    if ext in {".xlsx", ".xlsm"}:
        return parse_xlsx_file(file_bytes, filename)
    if ext == ".xls":
        return parse_xls_file(file_bytes, filename)
    if ext == ".pptx":
        return parse_pptx_file(file_bytes, filename)
    if ext in IMAGE_EXTENSIONS:
        raise ValueError("Image uploads need Gemini vision extraction before they can be used for RAG.")
    if ext in {".doc", ".ppt"}:
        raise ValueError(f"Legacy {ext} files need Gemini extraction or conversion to a modern Office format.")

    decoded = parse_text_file(file_bytes, filename)
    if decoded and _looks_like_text(decoded):
        return decoded
    raise ValueError(f"{ext or 'This file type'} is not supported for local text extraction.")


def _looks_like_text(value: str) -> bool:
    if not value.strip():
        return False
    sample = value[:2000]
    if "\x00" in sample:
        return False
    printable = sum(1 for char in sample if char.isprintable() or char.isspace())
    return printable / max(len(sample), 1) > 0.85


def supported_upload_extensions() -> list[str]:
    return sorted(
        {
            ".csv",
            ".doc",
            ".docx",
            ".htm",
            ".html",
            ".json",
            ".pdf",
            ".ppt",
            ".pptx",
            ".rtf",
            ".tsv",
            ".xls",
            ".xlsm",
            ".xlsx",
            *IMAGE_EXTENSIONS,
            *TEXT_LIKE_EXTENSIONS,
        }
    )


def chunk_text(
    text: str,
    *,
    chunk_size: int = 800,
    chunk_overlap: int = 150,
    source: str = "document",
    title: str = "Uploaded Document",
) -> list[DocumentChunk]:
    text = re.sub(r"\n{3,}", "\n\n", text.strip())
    if not text:
        return []

    paragraphs = re.split(r"\n\n+", text)
    chunks: list[DocumentChunk] = []
    current = ""
    chunk_idx = 0

    for paragraph in paragraphs:
        paragraph = paragraph.strip()
        if not paragraph:
            continue
        candidate = f"{current}\n\n{paragraph}".strip() if current else paragraph
        if len(candidate) > chunk_size and current:
            chunks.append(_build_chunk(source, title, current, chunk_idx))
            chunk_idx += 1
            overlap_text = current[-chunk_overlap:] if len(current) > chunk_overlap else current
            current = f"{overlap_text}\n\n{paragraph}".strip()
        else:
            current = candidate

    if current.strip():
        chunks.append(_build_chunk(source, title, current, chunk_idx))

    return chunks


def get_embeddings(
    texts: list[str],
    api_key: str | None = None,
    model: str = EMBEDDING_MODEL,
) -> list[list[float]]:
    cleaned_api_key = (api_key or "").strip()
    if cleaned_api_key:
        try:
            return _get_remote_embeddings(texts, cleaned_api_key, model=model)
        except Exception:
            pass
    return [_get_local_embedding(text).tolist() for text in texts]


def get_single_embedding(text: str, api_key: str | None = None) -> list[float]:
    results = get_embeddings([text], api_key=api_key)
    return results[0] if results else []


class SupabaseConfigError(RuntimeError):
    """Raised when the Supabase database connection is not configured."""


class SupabaseRestClient:
    def __init__(
        self,
        *,
        url: str | None = None,
        key: str | None = None,
        table_name: str | None = None,
        rpc_name: str | None = None,
    ) -> None:
        self.url = (url or os.getenv("SUPABASE_URL", "")).rstrip("/")
        self.key = (
            key
            or os.getenv("SUPABASE_SERVICE_ROLE_KEY", "")
            or os.getenv("SUPABASE_ANON_KEY", "")
        ).strip()
        self.table_name = (table_name or os.getenv("SUPABASE_DOCUMENT_TABLE", "document_chunks")).strip()
        self.rpc_name = (rpc_name or os.getenv("SUPABASE_MATCH_RPC", "match_document_chunks")).strip()

    @property
    def is_configured(self) -> bool:
        return bool(self.url and self.key)

    def upsert_chunks(self, chunks: list[DocumentChunk], *, session_id: str) -> None:
        rows = [
            {
                "id": chunk.id,
                "session_id": session_id,
                "source": chunk.source,
                "title": chunk.title,
                "content": chunk.text,
                "normalized_text": chunk.normalized_text,
                "tokens": sorted(chunk.tokens),
                "embedding": _format_pgvector(chunk.embedding),
            }
            for chunk in chunks
            if chunk.embedding
        ]
        if not rows:
            return
        self._request(
            "POST",
            f"/rest/v1/{self.table_name}",
            rows,
            extra_headers={"Prefer": "resolution=merge-duplicates,return=minimal"},
        )

    def delete_source(self, *, session_id: str, source: str) -> None:
        path = (
            f"/rest/v1/{self.table_name}"
            f"?session_id=eq.{_quote(session_id)}"
            f"&source=eq.{_quote(source)}"
        )
        self._request("DELETE", path, extra_headers={"Prefer": "return=minimal"})

    def match_chunks(
        self,
        query_embedding: list[float],
        *,
        session_id: str,
        top_k: int,
    ) -> list[DocumentChunk]:
        payload = {
            "query_embedding": _format_pgvector(query_embedding),
            "match_count": top_k,
            "match_session_id": session_id,
        }
        rows = self._request("POST", f"/rest/v1/rpc/{self.rpc_name}", payload)
        return [self._row_to_chunk(row) for row in rows if row.get("id")]

    def list_chunks(self, *, session_id: str, limit: int = 1000) -> list[DocumentChunk]:
        path = (
            f"/rest/v1/{self.table_name}"
            f"?session_id=eq.{_quote(session_id)}"
            "&select=id,source,title,content,normalized_text,tokens,embedding"
            f"&limit={limit}"
        )
        rows = self._request("GET", path)
        return [self._row_to_chunk(row) for row in rows if row.get("id")]

    def count_chunks(self, *, session_id: str) -> int:
        path = (
            f"/rest/v1/{self.table_name}"
            f"?session_id=eq.{_quote(session_id)}"
            "&select=id"
        )
        data, headers = self._request_raw(
            "GET",
            path,
            extra_headers={"Prefer": "count=exact"},
        )
        content_range = headers.get("Content-Range", "")
        if "/" in content_range:
            total = content_range.rsplit("/", 1)[-1]
            if total.isdigit():
                return int(total)
        return len(json.loads(data.decode("utf-8") or "[]"))

    def delete_session(self, *, session_id: str) -> None:
        path = f"/rest/v1/{self.table_name}?session_id=eq.{_quote(session_id)}"
        self._request("DELETE", path, extra_headers={"Prefer": "return=minimal"})

    def _row_to_chunk(self, row: dict[str, Any]) -> DocumentChunk:
        text = row.get("content") or row.get("text") or ""
        tokens = row.get("tokens") or []
        return DocumentChunk(
            id=str(row.get("id", "")),
            source=str(row.get("source", "document")),
            title=str(row.get("title", "Uploaded Document")),
            text=str(text),
            normalized_text=str(row.get("normalized_text") or _normalize_text(str(text))),
            tokens=set(tokens) if isinstance(tokens, list) else _tokenize(str(text)),
            embedding=_parse_embedding(row.get("embedding")),
        )

    def _request(
        self,
        method: str,
        path: str,
        payload: Any | None = None,
        *,
        extra_headers: dict[str, str] | None = None,
    ) -> Any:
        data, _ = self._request_raw(method, path, payload, extra_headers=extra_headers)
        if not data:
            return []
        return json.loads(data.decode("utf-8"))

    def _request_raw(
        self,
        method: str,
        path: str,
        payload: Any | None = None,
        *,
        extra_headers: dict[str, str] | None = None,
    ) -> tuple[bytes, dict[str, str]]:
        if not self.is_configured:
            raise SupabaseConfigError(
                "Supabase is not configured. Set SUPABASE_URL and SUPABASE_SERVICE_ROLE_KEY."
            )

        body = None if payload is None else json.dumps(payload).encode("utf-8")
        headers = {
            "apikey": self.key,
            "Authorization": f"Bearer {self.key}",
            "Content-Type": "application/json",
            "Accept": "application/json",
        }
        if extra_headers:
            headers.update(extra_headers)
        request = urllib.request.Request(
            url=f"{self.url}{path}",
            data=body,
            headers=headers,
            method=method,
        )
        try:
            with urllib.request.urlopen(request, timeout=30) as response:
                return response.read(), dict(response.headers.items())
        except urllib.error.HTTPError as exc:
            details = exc.read().decode("utf-8", errors="ignore")
            raise RuntimeError(f"Supabase HTTP error: {exc.code} {details}") from exc
        except urllib.error.URLError as exc:
            raise RuntimeError(f"Supabase connection error: {exc.reason}") from exc


class LocalMemoryChunkClient:
    """In-process fallback store for uploaded documents when Supabase is unavailable."""

    def __init__(self) -> None:
        self.rows: dict[str, DocumentChunk] = {}

    @property
    def is_configured(self) -> bool:
        return True

    def upsert_chunks(self, chunks: list[DocumentChunk], *, session_id: str) -> None:
        for chunk in chunks:
            self.rows[f"{session_id}:{chunk.id}"] = chunk

    def delete_source(self, *, session_id: str, source: str) -> None:
        self.rows = {
            key: chunk
            for key, chunk in self.rows.items()
            if key.split(":", 1)[0] != session_id or chunk.source != source
        }

    def match_chunks(
        self,
        query_embedding: list[float],
        *,
        session_id: str,
        top_k: int,
    ) -> list[DocumentChunk]:
        scored: list[tuple[float, DocumentChunk]] = []
        for key, chunk in self.rows.items():
            row_session_id, _ = key.split(":", 1)
            if row_session_id != session_id:
                continue
            scored.append((_cosine_similarity(query_embedding, chunk.embedding), chunk))
        scored.sort(key=lambda item: item[0], reverse=True)
        return [chunk for score, chunk in scored[:top_k] if score >= 0.0]

    def list_chunks(self, *, session_id: str, limit: int = 1000) -> list[DocumentChunk]:
        chunks = [
            chunk
            for key, chunk in self.rows.items()
            if key.split(":", 1)[0] == session_id
        ]
        return chunks[:limit]

    def count_chunks(self, *, session_id: str) -> int:
        return len(self.list_chunks(session_id=session_id))

    def delete_session(self, *, session_id: str) -> None:
        self.rows = {
            key: chunk
            for key, chunk in self.rows.items()
            if key.split(":", 1)[0] != session_id
        }


LOCAL_MEMORY_CHUNK_CLIENT = LocalMemoryChunkClient()


class SupabaseVectorStore:
    backend_name = "supabase"

    def __init__(
        self,
        *,
        session_id: str | None = None,
        client: Any | None = None,
    ) -> None:
        self.session_id = (session_id or os.getenv("SUPABASE_SESSION_ID", "default")).strip() or "default"
        self.client = client or SupabaseRestClient()
        self.backend_name = "local-memory" if isinstance(self.client, LocalMemoryChunkClient) else "supabase"
        self._cached_chunks: list[DocumentChunk] | None = None
        self._known_size: int | None = None
        self._last_scores: dict[str, float] = {}

    @property
    def size(self) -> int:
        if self._known_size is not None:
            return self._known_size
        self._known_size = self.client.count_chunks(session_id=self.session_id)
        return self._known_size

    @property
    def chunks(self) -> list[DocumentChunk]:
        if self._cached_chunks is None:
            self._cached_chunks = self.client.list_chunks(session_id=self.session_id)
            self._known_size = len(self._cached_chunks)
        return self._cached_chunks

    def add_chunks(self, chunks: list[DocumentChunk]) -> None:
        valid_chunks = [chunk for chunk in chunks if chunk.embedding]
        if not valid_chunks:
            return
        replaced_sources = {chunk.source for chunk in valid_chunks}
        if hasattr(self.client, "delete_source"):
            for source in replaced_sources:
                self.client.delete_source(session_id=self.session_id, source=source)
        self.client.upsert_chunks(valid_chunks, session_id=self.session_id)
        if self._cached_chunks is not None:
            existing = {
                chunk.id: chunk
                for chunk in self._cached_chunks
                if chunk.source not in replaced_sources
            }
            for chunk in valid_chunks:
                existing[chunk.id] = chunk
            self._cached_chunks = list(existing.values())
            self._known_size = len(self._cached_chunks)
        else:
            self._known_size = None

    def clear(self) -> None:
        self.client.delete_session(session_id=self.session_id)
        self._cached_chunks = []
        self._known_size = 0
        self._last_scores = {}

    def search(self, query_embedding: list[float], top_k: int = 100) -> list[RetrievalResult]:
        if not query_embedding:
            return []
        chunks = self.client.match_chunks(
            query_embedding,
            session_id=self.session_id,
            top_k=top_k,
        )
        self._last_scores = {
            chunk.id: max(0.0, min(1.0, _cosine_similarity(query_embedding, chunk.embedding)))
            for chunk in chunks
        }
        if self._cached_chunks is not None:
            existing = {chunk.id: chunk for chunk in self._cached_chunks}
            for chunk in chunks:
                existing[chunk.id] = chunk
            self._cached_chunks = list(existing.values())
        return [
            RetrievalResult(chunk=chunk, score=self._last_scores.get(chunk.id, 0.0))
            for chunk in chunks
            if self._last_scores.get(chunk.id, 0.0) >= 0.0
        ]


class RAGPipeline:
    """Project 2: uploaded-document assistant using Supabase retrieval with optional Gemini."""

    def __init__(
        self,
        api_key: str | None = None,
        *,
        gemini_client: Any | None = None,
        min_similarity: float = 0.0,
        session_id: str | None = None,
        vector_store: Any | None = None,
        supabase_client: Any | None = None,
    ) -> None:
        if api_key is None:
            configured_keys = [
                os.getenv("GEMINI_API_KEY", ""),
                *os.getenv("GEMINI_API_KEYS", "").split(","),
            ]
        else:
            configured_keys = [api_key]
        self.api_keys = list(
            dict.fromkeys(key.strip() for key in configured_keys if key.strip())
        )
        self.api_key = self.api_keys[0] if self.api_keys else ""
        self.gemini_client = gemini_client
        self.min_similarity = min_similarity
        self.vector_store = vector_store or self._build_vector_store(
            session_id=session_id,
            supabase_client=supabase_client,
        )
        self._processed_files: dict[str, str] = {}

    @property
    def has_documents(self) -> bool:
        return self.vector_store.size > 0

    @property
    def document_count(self) -> int:
        return len(self._processed_files)

    @property
    def chunk_count(self) -> int:
        return self.vector_store.size

    @property
    def vector_backend(self) -> str:
        return getattr(self.vector_store, "backend_name", "unknown")

    def ingest_file(self, file_bytes: bytes, filename: str) -> int:
        file_id = self._make_file_id(file_bytes, filename)
        if file_id in self._processed_files:
            return 0

        text = parse_uploaded_file(file_bytes, filename)
        return self.ingest_text(text, filename, file_id=file_id)

    def ingest_text(self, text: str, filename: str, *, file_id: str | None = None) -> int:
        file_id = file_id or self._make_text_id(text, filename)
        if file_id in self._processed_files:
            return 0
        if not text.strip():
            return 0

        chunks = chunk_text(
            text,
            source=filename,
            title=Path(filename).stem.replace("_", " ").replace("-", " ").title(),
        )
        if not chunks:
            return 0

        texts_to_embed = [chunk.text for chunk in chunks]
        embeddings = self._get_embeddings_with_fallback(texts_to_embed)
        for chunk, embedding in zip(chunks, embeddings):
            chunk.embedding = embedding

        self.vector_store.add_chunks(chunks)
        self._processed_files[file_id] = filename
        return len(chunks)

    def retrieve(self, query: str, top_k: int = 100) -> list[RetrievalResult]:
        if not self.has_documents:
            return []

        query_embedding = self._get_single_embedding_with_fallback(query)
        vector_results = self.vector_store.search(query_embedding, top_k=max(top_k * 3, 8))
        query_tokens = _tokenize(query)
        reranked: list[RetrievalResult] = []
        for result in vector_results:
            lexical_overlap = _token_overlap_score(query_tokens, result.chunk.tokens)
            final_score = (result.score * 0.74) + (lexical_overlap * 0.26)
            reranked.append(RetrievalResult(chunk=result.chunk, score=min(final_score, 1.0)))

        if not reranked:
            reranked = self._lexical_search(query, top_k=top_k)

        reranked.sort(key=lambda item: item.score, reverse=True)
        return self._dedupe_results(reranked, top_k=top_k)

    def answer_query(
        self,
        question: str,
        *,
        top_k: int = 24,
        chat_history: list[dict[str, Any]] | None = None,
        preferred_sources: list[str] | None = None,
        prefer_attached: bool = False,
    ) -> RAGResponse:
        clean_question = (question or "").strip()
        language = _detect_language(clean_question)
        contextualized_question = _contextualize_query(clean_question, chat_history)
        if not clean_question:
            return self._fallback_response(clean_question, language)

        scoped_to_uploaded_documents = (
            prefer_attached
            or bool(preferred_sources)
            or _looks_like_document_scope_question(clean_question)
        )

        if not self.has_documents:
            if scoped_to_uploaded_documents:
                return self._unreadable_documents_response(clean_question, language)
            freeform_response = self._freeform_response(contextualized_question, language, chat_history)
            if freeform_response:
                return freeform_response
            return self._no_documents_response(clean_question, language)

        retrieval_query = contextualized_question
        matches = self.retrieve(retrieval_query, top_k=top_k)
        if (not matches or matches[0].score < self.min_similarity) and (
            prefer_attached
            or preferred_sources
            or scoped_to_uploaded_documents
        ):
            attachment_matches = self._attachment_guided_matches(
                contextualized_question,
                preferred_sources=preferred_sources,
                top_k=top_k,
            )
            if attachment_matches:
                matches = attachment_matches
        if not matches or matches[0].score < self.min_similarity:
            if scoped_to_uploaded_documents:
                return self._fallback_response(clean_question, language)
            freeform_response = self._freeform_response(contextualized_question, language, chat_history)
            if freeform_response:
                return freeform_response
            return self._fallback_response(clean_question, language)

        matches_for_answer = matches[:8]
        source_urls = self._unique_sources(matches_for_answer)
        source_titles = self._unique_titles(matches_for_answer)
        if not self._gemini_configured():
            return self._model_required_documents_response(clean_question, language)

        try:
            answer_text = self._answer_with_gemini(
                user_question=contextualized_question,
                matches=matches_for_answer,
                language=language,
                chat_history=chat_history,
            )
            used_gemini = True
        except Exception:
            return self._model_required_documents_response(clean_question, language)

        return RAGResponse(
            question=clean_question,
            answer=answer_text,
            category="documents",
            source_url=source_urls[0] if source_urls else None,
            language=language,
            confidence=round(matches_for_answer[0].score, 3),
            found_match=True,
            used_gemini=used_gemini,
            matched_faq_id=matches_for_answer[0].chunk.id,
            source_urls=source_urls,
            source_titles=source_titles,
        )

    def _freeform_response(
        self,
        question: str,
        language: str,
        chat_history: list[dict[str, Any]] | None,
    ) -> RAGResponse | None:
        if not self.gemini_client or not getattr(self.gemini_client, "is_configured", False):
            return None
        try:
            answer = self.gemini_client.answer_freeform(
                user_question=question,
                language=language,
                chat_history=chat_history,
            )
        except Exception:
            return None
        return RAGResponse(
            question=question,
            answer=answer,
            category="chat",
            source_url=None,
            language=language,
            confidence=1.0,
            found_match=True,
            used_gemini=True,
            matched_faq_id=None,
            source_urls=[],
            source_titles=[],
        )

    def _model_required_documents_response(self, question: str, language: str) -> RAGResponse:
        answer = (
            "The live LLM is unavailable right now, so I cannot answer from uploaded documents at the moment. Please try again when the model is reachable."
            if language == "en"
            else "এই মুহূর্তে live LLM পাওয়া যাচ্ছে না, তাই আপলোড করা document থেকে উত্তর দিতে পারছি না। মডেল reachable হলে আবার চেষ্টা করুন।"
        )
        return RAGResponse(
            question=question,
            answer=answer,
            category="system",
            source_url=None,
            language=language,
            confidence=0.0,
            found_match=False,
            used_gemini=False,
            matched_faq_id=None,
            source_urls=[],
            source_titles=[],
        )

    def get_uploaded_files(self) -> list[str]:
        return sorted(set(self._processed_files.values()))

    def clear(self) -> None:
        self.vector_store.clear()
        self._processed_files.clear()

    def _lexical_search(self, query: str, *, top_k: int) -> list[RetrievalResult]:
        query_tokens = _tokenize(query)
        results: list[RetrievalResult] = []
        for chunk in self.vector_store.chunks:
            score = _token_overlap_score(query_tokens, chunk.tokens)
            if score > 0.08:
                results.append(RetrievalResult(chunk=chunk, score=score))
        results.sort(key=lambda item: item.score, reverse=True)
        return results[:top_k]

    def _build_vector_store(
        self,
        *,
        session_id: str | None,
        supabase_client: Any | None,
    ) -> SupabaseVectorStore:
        if supabase_client is not None:
            return SupabaseVectorStore(session_id=session_id, client=supabase_client)

        default_client = SupabaseRestClient()
        if getattr(default_client, "is_configured", False):
            return SupabaseVectorStore(session_id=session_id, client=default_client)

        return SupabaseVectorStore(session_id=session_id, client=LOCAL_MEMORY_CHUNK_CLIENT)

    def _attachment_guided_matches(
        self,
        question: str,
        *,
        preferred_sources: list[str] | None,
        top_k: int,
    ) -> list[RetrievalResult]:
        all_chunks = self.vector_store.chunks
        if not all_chunks:
            return []

        target_sources = self._target_sources_for_question(question, preferred_sources=preferred_sources)
        candidate_chunks = [chunk for chunk in all_chunks if not target_sources or chunk.source in target_sources]
        if not candidate_chunks:
            return []

        query_tokens = _tokenize(question)
        summary_request = _is_summary_request(question)
        document_scope_request = _looks_like_document_scope_question(question)
        ranked: list[RetrievalResult] = []
        for index, chunk in enumerate(candidate_chunks):
            overlap = _token_overlap_score(query_tokens, chunk.tokens)
            if preferred_sources and chunk.source in target_sources:
                overlap += 0.2
            if summary_request or not query_tokens:
                overlap = max(overlap, 0.28 - min(index * 0.01, 0.12))
            elif document_scope_request and chunk.source in target_sources:
                overlap = max(overlap, 0.18 - min(index * 0.01, 0.08))
            elif overlap <= 0:
                overlap = 0.12 if chunk.source in target_sources else 0.0
            if overlap > 0.1:
                ranked.append(RetrievalResult(chunk=chunk, score=min(overlap, 1.0)))

        ranked.sort(key=lambda item: item.score, reverse=True)
        return ranked[: max(3, min(top_k, 8))]

    def _target_sources_for_question(
        self,
        question: str,
        *,
        preferred_sources: list[str] | None,
    ) -> list[str]:
        if preferred_sources:
            preferred = {
                _normalize_source_name(source)
                for source in preferred_sources
                if _normalize_source_name(source)
            }
            matches = [
                chunk.source
                for chunk in self.vector_store.chunks
                if _normalize_source_name(chunk.source) in preferred
            ]
            if matches:
                return _unique_preserving_order(matches)

        normalized_question = _normalize_text(question)
        if normalized_question:
            mentioned = [
                chunk.source
                for chunk in self.vector_store.chunks
                if _normalize_source_name(chunk.source) and _normalize_source_name(chunk.source) in normalized_question
            ]
            if mentioned:
                return _unique_preserving_order(mentioned)

        recent_sources = list(dict.fromkeys(self._processed_files.values()))
        if recent_sources:
            return recent_sources[-1:]

        return [self.vector_store.chunks[-1].source]

    def _no_documents_response(self, question: str, language: str) -> RAGResponse:
        answer = (
            "Upload a document first, then ask me about it."
            if language == "en"
            else "আগে একটি document আপলোড করুন, তারপর সেটি নিয়ে প্রশ্ন করুন।"
        )
        return RAGResponse(
            question=question,
            answer=answer,
            category="system",
            source_url=None,
            language=language,
            confidence=0.0,
            found_match=False,
            used_gemini=False,
            matched_faq_id=None,
            source_urls=[],
            source_titles=[],
        )

    def _fallback_response(self, question: str, language: str) -> RAGResponse:
        answer = (
            "I couldn’t find that in the uploaded documents. Try a more specific question or upload a more relevant file."
            if language == "en"
            else "আপলোড করা document-এ এটা পাইনি। আরও নির্দিষ্ট প্রশ্ন করুন বা প্রাসঙ্গিক file আপলোড করুন।"
        )
        return RAGResponse(
            question=question,
            answer=answer,
            category="unknown",
            source_url=None,
            language=language,
            confidence=0.0,
            found_match=False,
            used_gemini=False,
            matched_faq_id=None,
            source_urls=[],
            source_titles=[],
        )

    def _unreadable_documents_response(self, question: str, language: str) -> RAGResponse:
        answer = (
            "I found the uploaded file, but I do not have readable text from it yet. Please upload the PDF or document again so I can extract and summarize its contents."
            if language == "en"
            else "আপলোড করা file পেয়েছি, কিন্তু এর readable text এখনও পাওয়া যায়নি। PDF বা document টি আবার upload করুন, অথবা text-select করে প্রশ্ন করুন।"
        )
        return RAGResponse(
            question=question,
            answer=answer,
            category="system",
            source_url=None,
            language=language,
            confidence=0.0,
            found_match=False,
            used_gemini=False,
            matched_faq_id=None,
            source_urls=[],
            source_titles=[],
        )

    def _answer_with_gemini(
        self,
        *,
        user_question: str,
        matches: list[RetrievalResult],
        language: str,
        chat_history: list[dict[str, Any]] | None = None,
    ) -> str:
        if self.gemini_client and hasattr(self.gemini_client, "answer_from_context"):
            return self.gemini_client.answer_from_context(
                user_question=user_question,
                matches=matches,
                language=language,
                chat_history=chat_history,
            )

        api_key = self._gemini_api_key()
        if not api_key:
            raise RuntimeError("Gemini API key is not configured.")

        model = self._gemini_model()
        prompt = self._build_prompt(
            user_question=user_question,
            matches=matches,
            language=language,
            chat_history=chat_history,
        )
        payload = {
            "contents": [{"parts": [{"text": prompt}]}],
            "generationConfig": {
                "temperature": 0.4,
                "topP": 0.9,
                "maxOutputTokens": int(os.getenv("GEMINI_MAX_OUTPUT_TOKENS", "8192")),
            },
        }
        url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent"
        request = urllib.request.Request(
            url=url,
            data=json.dumps(payload).encode("utf-8"),
            headers={
                "Content-Type": "application/json",
                "x-goog-api-key": api_key,
            },
            method="POST",
        )

        try:
            with urllib.request.urlopen(request, timeout=30) as response:
                data = json.loads(response.read().decode("utf-8"))
        except urllib.error.HTTPError as exc:
            details = exc.read().decode("utf-8", errors="ignore")
            raise RuntimeError(f"Gemini API HTTP error: {exc.code} {details}") from exc
        except urllib.error.URLError as exc:
            raise RuntimeError(f"Gemini API connection error: {exc.reason}") from exc

        text = self._extract_model_text(data)
        if not text:
            raise RuntimeError("Gemini API returned an empty response.")
        return self._clean_model_text(text)

    def _build_prompt(
        self,
        *,
        user_question: str,
        matches: list[RetrievalResult],
        language: str,
        chat_history: list[dict[str, Any]] | None = None,
    ) -> str:
        preferred_language = "Bangla" if language == "bn" else "English"
        formatted_history = _format_history(chat_history)
        context_blocks: list[str] = []
        for index, match in enumerate(matches, start=1):
            context_blocks.append(
                f"[Document {index}]\n"
                f"Title: {match.chunk.title}\n"
                f"Source: {match.chunk.source}\n"
                f"Excerpt: {match.chunk.text}\n"
            )
        joined_context = "\n".join(context_blocks)
        return (
            "You are the official Daffodil International University (DIU) AI Assistant for uploaded university documents.\n"
            "Answer directly and naturally, but keep the tone mature, polished, and professionally organized.\n"
            "Do not introduce yourself as Gemini, Google, a language model, or a generic AI model.\n"
            "Use uploaded document excerpts as helpful context, not as a restriction.\n"
            "If the documents are incomplete, still give the best useful answer and avoid sounding blocked by missing context.\n"
            "If the user asks for official DIU facts that are not confirmed by the uploaded material, say that you do not have the official documentation for that point.\n"
            "If the question is clearly unrelated to DIU, academia, student life, or the uploaded university documents, politely decline and redirect to DIU matters.\n"
            "Use the user's preferred clean assistant format: direct overview first, then clear numbered sections, bullets, short explanatory paragraphs, and Markdown tables where useful.\n"
            "For broad document questions, include the main categories, subpoints, rules, values, exceptions, and notes found in the uploaded excerpts. Prefer organized completeness over a short summary.\n"
            "Avoid abrupt bullet-only answers; make the reasoning, categories, conditions, and verification points easy to scan and understand.\n"
            "Close with a short natural takeaway only when it adds clarity; do not force a labeled **Conclusion:** line or generic **Next Steps:** section. Do not add a follow-up question after it.\n"
            f"Preferred response language: {preferred_language}\n"
            "Recent session context (use only when it helps resolve the current follow-up):\n"
            f"{formatted_history}\n\n"
            "Use session history only if the current question is clearly referring back to it. Ignore stale history for greetings or topic changes.\n\n"
            f"User question: {user_question}\n\n"
            f"Optional uploaded document context:\n{joined_context}\n"
        )

    def _extract_model_text(self, payload: dict[str, Any]) -> str:
        parts: list[str] = []
        for candidate in payload.get("candidates", []):
            content = candidate.get("content", {})
            for part in content.get("parts", []):
                text = part.get("text")
                if text:
                    parts.append(text.strip())
        return "\n".join(part for part in parts if part).strip()

    def _clean_model_text(self, text: str) -> str:
        text = re.sub(r"^\s*answer\s*:\s*", "", text.strip(), flags=re.IGNORECASE)
        text = self._sanitize_identity_leaks(text)
        text = re.sub(r"\n{3,}", "\n\n", text)
        return text.strip()

    def _sanitize_identity_leaks(self, text: str) -> str:
        replacements = (
            (
                r"^\s*(?:i am|i'm)\s+(?:google(?:'s)?\s+)?gemini(?:\s+\w+)*[,.:\-\s]*",
                "I’m the DIU Assistant. ",
            ),
            (
                r"^\s*(?:i am|i'm)\s+(?:an?\s+)?(?:ai language model|ai assistant|ai model|language model|large language model|ai)[,.:\-\s]*",
                "I’m the DIU Assistant. ",
            ),
            (
                r"^\s*as\s+(?:google(?:'s)?\s+)?gemini(?:\s+\w+)*[,:]?\s*",
                "",
            ),
            (
                r"^\s*as\s+(?:an?\s+)?(?:ai language model|ai assistant|ai model|language model|large language model|ai)[,:]?\s*",
                "",
            ),
        )
        cleaned = text.strip()
        for pattern, replacement in replacements:
            cleaned = re.sub(pattern, replacement, cleaned, flags=re.IGNORECASE)
        cleaned = re.sub(
            r"^\s*(?:i am|i'm)\s+(?:google(?:'s)?\s+)?gemini(?:\s+\w+)*[,.:\-\s]*",
            "I’m the DIU Assistant. ",
            cleaned,
            flags=re.IGNORECASE,
        )
        cleaned = re.sub(
            r"\b(?:Google\s+Gemini|Gemini)\b(?=\s+(?:assistant|model|api)\b)",
            "the DIU Assistant",
            cleaned,
            flags=re.IGNORECASE,
        )
        cleaned = re.sub(
            r"\b(?:powered|built|run)\s+by\s+(?:Google\s+Gemini|Gemini)\b",
            "",
            cleaned,
            flags=re.IGNORECASE,
        )
        cleaned = re.sub(
            r"\bmy capabilities\s+(?:are\s+)?(?:powered|built|run)\s+by\s+(?:Google\s+Gemini|Gemini)\b[,.:\-\s]*",
            "",
            cleaned,
            flags=re.IGNORECASE,
        )
        cleaned = re.sub(
            r"\b(?:the\s+)?google(?:'s)?\s+(?:ai language model|ai assistant|ai model|language model|large language model)\b[,.:\-\s]*",
            "",
            cleaned,
            flags=re.IGNORECASE,
        )
        cleaned = re.sub(r"[^\S\n]{2,}", " ", cleaned)
        cleaned = re.sub(r"[^\S\n]+([,.;:])", r"\1", cleaned)
        cleaned = re.sub(
            r"\b(?:as\s+)?an?\s+(?:AI language model|AI assistant|AI model|language model|large language model|AI)\b",
            "as the DIU Assistant",
            cleaned,
            flags=re.IGNORECASE,
        )
        cleaned = re.sub(r"^(I’m the DIU Assistant\.\s*){2,}", "I’m the DIU Assistant. ", cleaned)
        return cleaned.strip()

    def _gemini_configured(self) -> bool:
        if self.gemini_client and hasattr(self.gemini_client, "is_configured"):
            return bool(getattr(self.gemini_client, "is_configured"))
        return bool(self._gemini_api_key())

    def _gemini_api_key(self) -> str:
        if self.gemini_client and getattr(self.gemini_client, "api_key", ""):
            return getattr(self.gemini_client, "api_key", "").strip()
        return self.api_key

    def _get_embeddings_with_fallback(self, texts: list[str]) -> list[list[float]]:
        for api_key in self.api_keys:
            try:
                return get_embeddings(texts, api_key=api_key)
            except Exception:
                continue
        return get_embeddings(texts, api_key=None)

    def _get_single_embedding_with_fallback(self, text: str) -> list[float]:
        results = self._get_embeddings_with_fallback([text])
        return results[0] if results else []

    def _gemini_model(self) -> str:
        if self.gemini_client and getattr(self.gemini_client, "model", ""):
            return getattr(self.gemini_client, "model")
        return os.getenv("GEMINI_MODEL", "gemini-3-flash-preview")

    def _unique_sources(self, matches: list[RetrievalResult]) -> list[str]:
        seen: set[str] = set()
        sources: list[str] = []
        for match in matches:
            if match.chunk.source not in seen:
                seen.add(match.chunk.source)
                sources.append(match.chunk.source)
        return sources

    def _unique_titles(self, matches: list[RetrievalResult]) -> list[str]:
        seen: set[str] = set()
        titles: list[str] = []
        for match in matches:
            if match.chunk.source not in seen:
                seen.add(match.chunk.source)
                titles.append(match.chunk.title)
        return titles

    def _dedupe_results(self, matches: list[RetrievalResult], *, top_k: int) -> list[RetrievalResult]:
        chosen: list[RetrievalResult] = []
        seen_sources: set[str] = set()
        for match in matches:
            if match.chunk.source not in seen_sources:
                chosen.append(match)
                seen_sources.add(match.chunk.source)
            if len(chosen) >= top_k:
                break
        return chosen

    def _make_file_id(self, file_bytes: bytes, filename: str) -> str:
        digest = hashlib.sha1(file_bytes).hexdigest()[:16]
        return f"{filename}:{digest}"

    def _make_text_id(self, text: str, filename: str) -> str:
        digest = hashlib.sha1(text.encode("utf-8", errors="ignore")).hexdigest()[:16]
        return f"{filename}:{digest}"


def _build_chunk(source: str, title: str, text: str, chunk_idx: int) -> DocumentChunk:
    normalized_text = _normalize_text(text)
    return DocumentChunk(
        id=f"{_make_id(source)}-{chunk_idx}",
        source=source,
        title=title,
        text=text.strip(),
        normalized_text=normalized_text,
        tokens=_tokenize(text),
    )


def _normalize_text(text: str) -> str:
    text = text.lower()
    text = re.sub(r"[^0-9a-z\u0980-\u09ff\s]", " ", text)
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def _tokenize(text: str) -> set[str]:
    normalized = _normalize_text(text)
    return {
        token
        for token in normalized.split()
        if token and token not in STOPWORDS and len(token) > 1
    }


def _token_overlap_score(query_tokens: set[str], chunk_tokens: set[str]) -> float:
    if not query_tokens or not chunk_tokens:
        return 0.0
    return len(query_tokens & chunk_tokens) / max(min(len(query_tokens), 6), 1)


def _detect_language(text: str) -> str:
    bangla_chars = len(re.findall(r"[\u0980-\u09FF]", text))
    english_chars = len(re.findall(r"[A-Za-z]", text))
    if bangla_chars >= 2 and bangla_chars >= max(1, english_chars * 0.25):
        return "bn"
    return "en"


def _trim_to_sentences(text: str, *, limit: int = 2) -> str:
    sentences = re.split(r"(?<=[.!?])\s+", text.strip())
    selected = [sentence.strip() for sentence in sentences if sentence.strip()][:limit]
    result = " ".join(selected)
    return result[:520].strip()


def _clean_sentence(sentence: str) -> str:
    sentence = re.sub(r"\s+", " ", sentence).strip(" -|")
    if len(sentence) < 30:
        return ""
    noisy = (
        "subscribe",
        "copyright",
        "all rights reserved",
        "click here",
        "read more",
        "login",
    )
    if any(term in sentence.lower() for term in noisy):
        return ""
    return sentence[:360].strip()


def _contextualize_query(
    question: str,
    chat_history: list[dict[str, Any]] | None,
) -> str:
    clean_question = str(question or "").strip()
    if not clean_question or not chat_history or not _looks_like_followup(clean_question):
        return clean_question

    topic_index = _latest_substantive_user_topic(chat_history)
    scoped_history = chat_history[topic_index:] if topic_index >= 0 else chat_history
    session_topic = ""
    if topic_index >= 0:
        session_topic = _history_excerpt_for_retrieval(str(chat_history[topic_index].get("content", "")).strip())

    recent_context = _recent_followup_context(clean_question, scoped_history)
    if not session_topic and not recent_context:
        return clean_question

    return (
        f"{clean_question}\n\n"
        f"Session topic: {session_topic or recent_context}\n"
        f"Relevant recent context: {recent_context}"
    )


def _latest_substantive_user_topic(chat_history: list[dict[str, Any]]) -> int:
    for index in range(len(chat_history) - 1, -1, -1):
        message = chat_history[index]
        if message.get("role") != "user":
            continue
        content = str(message.get("content", "")).strip()
        if content and not _looks_like_followup(content):
            return index

    for index in range(len(chat_history) - 1, -1, -1):
        message = chat_history[index]
        if message.get("role") == "user" and str(message.get("content", "")).strip():
            return index

    return -1


def _recent_followup_context(
    question: str,
    chat_history: list[dict[str, Any]],
) -> str:
    normalized_question = _normalize_text(question)
    excerpts: list[str] = []
    for message in chat_history[-8:]:
        role = message.get("role")
        if role not in {"user", "assistant"}:
            continue
        content = str(message.get("content", "")).strip()
        if not content or _normalize_text(content) == normalized_question:
            continue
        excerpt = _history_excerpt_for_retrieval(content)
        if excerpt:
            excerpts.append(excerpt)
    return " ".join(excerpts[-4:])


def _history_excerpt_for_retrieval(content: str) -> str:
    cleaned = re.sub(r"`([^`]+)`", r"\1", content)
    cleaned = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", cleaned)
    cleaned = re.sub(r"[*_#>|-]+", " ", cleaned)
    cleaned = re.sub(r"\s+", " ", cleaned).strip()
    return cleaned[:320]


def _format_history(chat_history: list[dict[str, Any]] | None) -> str:
    if not chat_history:
        return "(none)"
    lines: list[str] = []
    for index, message in enumerate(chat_history[-12:], start=1):
        role = message.get("role", "user")
        mode = str(message.get("mode", "")).strip()
        content = str(message.get("content", "")).strip()
        if content:
            label = f"{role} [{mode}]" if mode else str(role)
            lines.append(f"{index}. {label}: {content[:1000]}")
    return "\n".join(lines) if lines else "(none)"


def _looks_like_followup(question: str) -> bool:
    normalized = _normalize_text(question)
    tokens = normalized.split()
    if not normalized:
        return False
    vague_followups = (
        "more",
        "more detail",
        "more details",
        "more detailed",
        "details",
        "detail please",
        "explain more",
        "tell me more",
        "elaborate",
        "chart",
        "table",
        "summarize it",
        "summarise it",
    )
    if len(tokens) <= 4 and any(phrase in normalized for phrase in vague_followups):
        return True
    if len(tokens) <= 2:
        return any(
            token in {"it", "this", "that", "them", "those", "these", "same", "related"}
            for token in tokens
        )
    if len(tokens) <= 5 and any(
        token in {"it", "this", "that", "them", "those", "these", "same", "related"}
        for token in tokens
    ):
        return True
    starters = (
        "what about",
        "how about",
        "and",
        "also",
        "then",
        "that",
        "this",
        "it",
        "them",
        "those",
        "these",
        "same",
        "related",
        "for this",
        "for that",
        "in this",
        "in that",
        "আর",
        "তাহলে",
    )
    if normalized.startswith(starters):
        return True
    return len(tokens) <= 10 and any(
        token in {"it", "this", "that", "them", "those", "these", "same", "related"}
        for token in tokens
    )


def _intent_terms_for_question(question: str) -> set[str]:
    lowered = _normalize_text(question)
    groups = [
        {"admission", "apply", "eligibility", "requirement", "requirements", "ভর্তি"},
        {"scholarship", "waiver", "financial", "aid", "স্কলারশিপ", "ওয়েভার"},
        {"course", "program", "curriculum", "department", "কোর্স", "প্রোগ্রাম"},
        {"fee", "fees", "tuition", "cost", "payment", "ফি"},
    ]
    required: set[str] = set()
    for group in groups:
        if any(term in lowered for term in group):
            required.update(group)
    return required


def _is_summary_request(question: str) -> bool:
    normalized = _normalize_text(question)
    return bool(
        re.search(
            r"\b(summarize|summarise|summary|overview|main points?|key points?|what is in|what does this say)\b",
            normalized,
        )
    )


def _looks_like_document_scope_question(question: str) -> bool:
    normalized = _normalize_text(question)
    return bool(
        re.search(
            r"\b(attached|attachment|upload|uploaded|document|documents|file|files|pdf|image|screenshot|slide|spreadsheet)\b",
            normalized,
        )
    ) or _is_summary_request(question)


def _normalize_source_name(source: str) -> str:
    basename = Path(source or "").name
    stem = Path(basename).stem
    normalized = _normalize_text(f"{basename} {stem}")
    return normalized.strip()


def _unique_preserving_order(values: list[str]) -> list[str]:
    return list(dict.fromkeys(values))


def _get_local_embedding(text: str, dimension: int = EMBEDDING_DIMENSION) -> np.ndarray:
    vector = np.zeros(dimension, dtype=np.float32)
    tokens = list(_tokenize(text))
    if not tokens:
        tokens = re.findall(r"[0-9a-z\u0980-\u09ff]+", _normalize_text(text))

    for token in tokens:
        digest = hashlib.sha256(token.encode("utf-8")).digest()
        index = int.from_bytes(digest[:2], "big") % dimension
        sign = 1.0 if digest[2] % 2 == 0 else -1.0
        weight = 1.0 + (min(len(token), 10) / 10.0)
        vector[index] += sign * weight

    norm = np.linalg.norm(vector)
    if norm == 0:
        return vector
    return vector / norm


def _get_remote_embeddings(
    texts: list[str],
    api_key: str,
    *,
    model: str = EMBEDDING_MODEL,
) -> list[list[float]]:
    url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:batchEmbedContents"
    requests_body = [
        {"model": f"models/{model}", "content": {"parts": [{"text": text[:2048]}]}}
        for text in texts
    ]

    payload = json.dumps({"requests": requests_body}).encode("utf-8")
    request = urllib.request.Request(
        url=url,
        data=payload,
        headers={
            "Content-Type": "application/json",
            "x-goog-api-key": api_key,
        },
        method="POST",
    )

    with urllib.request.urlopen(request, timeout=30) as response:
        data = json.loads(response.read().decode("utf-8"))

    embeddings: list[list[float]] = []
    for embedding in data.get("embeddings", []):
        embeddings.append(embedding.get("values", []))
    return embeddings


def _make_id(source: str) -> str:
    cleaned = re.sub(r"[^a-z0-9]", "-", source.lower())
    cleaned = re.sub(r"-+", "-", cleaned).strip("-")
    return cleaned[:40] or "doc"


def _format_pgvector(values: list[float]) -> str:
    cleaned = [float(value) for value in values]
    return "[" + ",".join(f"{value:.8f}" for value in cleaned) + "]"


def _parse_embedding(value: Any) -> list[float]:
    if isinstance(value, list):
        return [float(item) for item in value]
    if isinstance(value, str):
        stripped = value.strip().strip("[]")
        if not stripped:
            return []
        return [float(item) for item in stripped.split(",") if item.strip()]
    return []


def _cosine_similarity(left: list[float], right: list[float]) -> float:
    if not left or not right:
        return 0.0
    left_vec = np.array(left, dtype=np.float32)
    right_vec = np.array(right, dtype=np.float32)
    if left_vec.size != right_vec.size:
        return 0.0
    left_norm = np.linalg.norm(left_vec)
    right_norm = np.linalg.norm(right_vec)
    if left_norm == 0 or right_norm == 0:
        return 0.0
    return float(np.dot(left_vec / left_norm, right_vec / right_norm))


def _quote(value: str) -> str:
    return urllib.parse.quote(value, safe="")
